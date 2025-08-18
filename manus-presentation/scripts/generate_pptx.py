from pptx import Presentation
from pptx.util import Inches, Pt
import sys
import os

def markdown_to_slides(md_path):
    with open(md_path, 'r', encoding='utf-8') as f:
        content = f.read()
    parts = [p.strip() for p in content.split('---') if p.strip()]
    slides = []
    for p in parts:
        lines = [l.strip() for l in p.splitlines() if l.strip()]
        if not lines:
            continue
        title = lines[0].lstrip('#').strip()
        body_lines = []
        for l in lines[1:]:
            if l.startswith('#'):
                body_lines.append(l.lstrip('#').strip())
            else:
                body_lines.append(l)
        slides.append((title, body_lines))
    return slides


def create_pptx(slides, out_path):
    prs = Presentation()
    # set default font size for body
    for i, (title, body) in enumerate(slides):
        if i == 0:
            # title slide layout
            slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(slide_layout)
            if title:
                slide.shapes.title.text = title
            if body:
                # put subtitle if available
                try:
                    slide.placeholders[1].text = '\n'.join(body)
                except Exception:
                    pass
        else:
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            if title:
                slide.shapes.title.text = title
            if body:
                body_shape = slide.shapes.placeholders[1]
                tf = body_shape.text_frame
                for idx, line in enumerate(body):
                    if idx == 0:
                        tf.text = line
                    else:
                        p = tf.add_paragraph()
                        p.text = line
                        p.level = 0
                        p.font.size = Pt(14)
    # ensure output dir exists
    os.makedirs(os.path.dirname(out_path), exist_ok=True)
    prs.save(out_path)


if __name__ == '__main__':
    if len(sys.argv) < 3:
        print('Usage: python generate_pptx.py <slides.md> <output.pptx>')
        sys.exit(1)
    md = sys.argv[1]
    out = sys.argv[2]
    slides = markdown_to_slides(md)
    create_pptx(slides, out)
    print(f'Generated PPTX: {out}')
